"""Builds presentation/presentation.pptx -- calm dark slate-blue deck, 9 slides,
16:9. Run AFTER metrics.csv + qini curves + business evaluation all exist.

Design system: dark background, muted steel-blue accent, uppercase letter-spaced
kickers (tagged to the AI Talent Hub evaluation dimension each slide speaks
to -- Development & Engineering / Data Science / AI Application / Product
Thinking / Motivation, per docs/itmo_requirements.md), card grids with a
colored left accent bar, and a consistent footer breadcrumb + page counter.
"""
import json
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

sys.path.insert(0, os.path.dirname(__file__))
from theme import (  # noqa: E402
    BG, BG_ALT, CARD_BG, CARD_BORDER, PRIMARY, PRIMARY_DARK, AMBER, GOLD, GREEN,
    GRAY_TEXT, WHITE, MODEL_COLORS, MODEL_LABELS, FONT,
)

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.55
TOTAL_SLIDES = 9


def rgb(hexstr):
    h = hexstr.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


prs = Presentation()
prs.slide_width = Inches(SLIDE_W)
prs.slide_height = Inches(SLIDE_H)
BLANK = prs.slide_layouts[6]


def new_slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = rgb(BG)
    return s


def add_rect(slide, x, y, w, h, color, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(color)
    if line:
        shp.line.color.rgb = rgb(CARD_BORDER)
        shp.line.width = Pt(0.75)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_rounded_rect(slide, x, y, w, h, fill_color, border_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = 0.06
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill_color)
    if border_color:
        shp.line.color.rgb = rgb(border_color)
        shp.line.width = Pt(0.75)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size, color=WHITE, bold=False, italic=False,
             align=PP_ALIGN.LEFT, anchor=None, line_spacing=1.0, font=FONT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font
        run.font.color.rgb = rgb(color)
    return tb


def add_kicker(slide, text, x=MARGIN, y=0.42, color=PRIMARY):
    add_rect(slide, x, y + 0.06, 0.32, 0.05, color)
    spaced = " ".join(list(text.upper()))
    add_text(slide, x + 0.45, y - 0.08, SLIDE_W - 2 * x, 0.35, spaced, 11.5, color=color, bold=True)


def add_title(slide, text, x=MARGIN, y=0.78, size=30, w=None):
    add_text(slide, x, y, (w or SLIDE_W - 2 * x), 1.0, text, size, color=WHITE, bold=True, line_spacing=1.03)


def add_subtitle(slide, text, x=MARGIN, y=1.55, size=14.5, w=None):
    add_text(slide, x, y, (w or SLIDE_W - 2 * x), 0.6, text, size, color=GRAY_TEXT)


def add_footer(slide, page):
    add_rect(slide, 0, SLIDE_H - 0.5, SLIDE_W, 0.012, CARD_BORDER)
    add_text(slide, MARGIN, SLIDE_H - 0.42, 8.0, 0.3,
             "Uplift Modeling · Junior ML Contest · AI Talent Hub", 10.5, color=GRAY_TEXT)
    add_text(slide, SLIDE_W - MARGIN - 1.5, SLIDE_H - 0.42, 1.5, 0.3,
             f"{page:02d} / {TOTAL_SLIDES:02d}", 10.5, color=GRAY_TEXT, align=PP_ALIGN.RIGHT)


def _estimate_wrapped_lines(text, width_in, font_size, bold=True):
    # Rough heuristic: bold Calibri averages ~0.55*font_size(pt) px width per
    # char at 72dpi -- good enough to decide 1 vs 2 wrapped lines for layout,
    # not for exact typesetting.
    avg_char_width_in = (font_size * (0.62 if bold else 0.55)) / 72.0
    chars_per_line = max(int(width_in / avg_char_width_in), 1)
    return max(1, -(-len(text) // chars_per_line))  # ceil division


def add_card(slide, x, y, w, h, title, body, accent=PRIMARY, title_size=15, body_size=12.5):
    add_rounded_rect(slide, x, y, w, h, CARD_BG, CARD_BORDER)
    add_rect(slide, x, y + 0.12, 0.06, h - 0.24, accent)
    pad = 0.28
    title_lines = _estimate_wrapped_lines(title, w - 2 * pad, title_size)
    title_box_h = title_lines * (title_size / 72 * 1.25) + 0.06
    add_text(slide, x + pad, y + 0.16, w - 2 * pad, title_box_h + 0.15, title, title_size, color=WHITE, bold=True, line_spacing=1.05)
    add_text(slide, x + pad, y + 0.16 + title_box_h + 0.12, w - 2 * pad, h - 0.6, body, body_size, color=GRAY_TEXT, line_spacing=1.08)


def add_stat_tile(slide, x, y, w, h, number, label, desc, number_color=PRIMARY, number_size=40):
    add_rounded_rect(slide, x, y, w, h, CARD_BG, CARD_BORDER)
    add_rect(slide, x, y + 0.1, 0.06, h - 0.2, number_color)
    pad = 0.26
    add_text(slide, x + pad, y + 0.12, w - 2 * pad, 0.75, number, number_size, color=number_color, bold=True)
    add_text(slide, x + pad, y + 0.12 + number_size / 72 + 0.05, w - 2 * pad, 0.4, label, 14, color=WHITE, bold=True)
    add_text(slide, x + pad, y + h - 0.62, w - 2 * pad, 0.55, desc, 10.5, color=GRAY_TEXT, line_spacing=1.05)


def add_image_fit(slide, path, x, y, max_w, max_h):
    im = Image.open(path)
    ar = im.width / im.height
    w, h = max_w, max_w / ar
    if h > max_h:
        h = max_h
        w = max_h * ar
    left = x + (max_w - w) / 2
    top = y + (max_h - h) / 2
    slide.shapes.add_picture(path, Inches(left), Inches(top), width=Inches(w), height=Inches(h))


with open("artifacts/metrics.json") as f:
    metrics = json.load(f)
metrics_by_model = {m["model"]: m for m in metrics}

page = 1  # slide 1 (title) carries no footer/number, so the first numbered slide is #2


def next_page():
    global page
    page += 1
    return page


# ======================================================================
# Slide 1 — Title
# ======================================================================
s = new_slide()
add_rect(s, SLIDE_W * 0.62, 0, SLIDE_W * 0.38, SLIDE_H, BG_ALT)

# decorative ascending "uplift" bars on the right panel
bar_x0 = SLIDE_W * 0.68
bar_colors = ["#243447", "#2C4A61", "#355F79", PRIMARY, GOLD, GREEN]
bar_heights = [0.9, 1.5, 2.1, 2.8, 3.6, 4.5]
bar_w = 0.5
gap = 0.16
for i, (hgt, col) in enumerate(zip(bar_heights, bar_colors)):
    bx = bar_x0 + i * (bar_w + gap)
    by = 5.6 - hgt
    add_rect(s, bx, by, bar_w, hgt, col)
add_text(s, bar_x0, 5.75, 4.0, 0.35, "иллюстрация: прирост эффекта коммуникации", 10.5, color=GRAY_TEXT)

add_text(s, MARGIN, 0.5, 6.5, 0.35, "X 5   R E T A I L H E R O   ·   J U N I O R   M L   C O N T E S T", 12,
         color=PRIMARY, bold=True)
add_text(s, MARGIN, 1.9, 7.0, 2.0,
         "Uplift Modeling\nс последовательностью\nповедения клиента", 38, color=WHITE, bold=True, line_spacing=1.05)
add_text(s, MARGIN, 4.05, 7.0, 0.8,
         "Оценка приростного эффекта клиентских коммуникаций: классические\nCatBoost-learner'ы против Transformer на сырых чеках", 15,
         color=GRAY_TEXT, line_spacing=1.15)

add_rect(s, MARGIN, 5.35, 0.4, 0.03, PRIMARY)
add_text(s, MARGIN, 5.5, 6.5, 0.4, "Галиб Байрамов", 18, color=WHITE, bold=True)
add_text(s, MARGIN, 5.9, 6.7, 0.35, "Продуктовый аналитик, Т-Банк · маркетинговая аналитика, привлечение B2B", 12.5, color=GRAY_TEXT)
add_text(s, MARGIN, 6.25, 6.7, 0.35, "ИТМО AI Talent Hub · Junior ML Contest · 2026", 12.5, color=GRAY_TEXT)

# ======================================================================
# Slide 2 — Motivation (per the user's exact brief)
# ======================================================================
s = new_slide()
add_kicker(s, "Мотивация исследования", color=GOLD)
add_title(s, "Почему я выбрал эту задачу")
add_subtitle(s, "Профессиональный опыт → исследовательский вопрос → самостоятельная проверка на открытых данных")

col_w = (SLIDE_W - 2 * MARGIN - 2 * 0.35) / 3
col_y = 2.25
col_h = 4.15

add_card(
    s, MARGIN, col_y, col_w, col_h,
    "Профессиональный контекст",
    "Product Analyst, Т-Банк\nКоманда роста маркетинговой аналитики\nПривлечение B2B-клиентов\n\n"
    "— оценка маркетинговых запусков\n— сегментация клиентов\n— A/B-тесты\n— продуктовые метрики",
    accent=PRIMARY_DARK, body_size=13,
)
add_card(
    s, MARGIN + col_w + 0.35, col_y, col_w, col_h,
    "Возникший вопрос",
    "Средний эффект A/B-теста показывает, работает ли коммуникация в целом, "
    "но не отвечает, какому клиенту она помогла.\n\n"
    "Можно ли оценить индивидуальный эффект коммуникации?",
    accent=AMBER, body_size=13.5,
)
add_card(
    s, MARGIN + 2 * (col_w + 0.35), col_y, col_w, col_h,
    "Самостоятельное исследование",
    "Открытые данные X5 RetailHero\nUplift-моделирование\nT-Learner, X-Learner, Transformer\n\n"
    "Профессиональная задача стала мотивацией для самостоятельного ML-исследования.",
    accent=GREEN, body_size=13,
)

for i in range(2):
    ax = MARGIN + col_w + i * (col_w + 0.35) - 0.28
    arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(ax), Inches(col_y + col_h / 2 - 0.1), Inches(0.22), Inches(0.2))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = rgb(GRAY_TEXT); arrow.line.fill.background(); arrow.shadow.inherit = False

add_text(s, MARGIN, col_y + col_h + 0.25, SLIDE_W - 2 * MARGIN, 0.5,
          "Проект выполнен на открытых данных X5 RetailHero. Данные, код и внутренние результаты Т-Банка не использовались.",
          12.5, color=GRAY_TEXT, italic=True)
add_footer(s, next_page())

# ======================================================================
# Slide 3 — Business problem
# ======================================================================
s = new_slide()
add_kicker(s, "Product Thinking", color=PRIMARY)
add_title(s, "Бизнес-проблема")
add_subtitle(s, "Коммуникация стоит денег — обращаться нужно к тем, кто изменит поведение именно благодаря ей")

cards = [
    ("Стоимость и риск", "Отправка SMS стоит денег и может раздражать клиента, поэтому обращаться нужно не ко всем, а к тем, кому это действительно поможет.", PRIMARY),
    ("Отклик — это не прирост", "Модель отклика находит тех, кто и так с высокой вероятностью купит. Из-за этого она выбирает не тех, кому нужна коммуникация, а тех, кто обошёлся бы и без неё.", AMBER),
    ("Причинная оценка", "Значит, вместо вероятности покупки нужно оценивать, изменится ли поведение клиента именно благодаря коммуникации.", GOLD),
    ("Цель", "Тогда при фиксированном бюджете рассылки можно получить максимальный прирост конверсий, а не просто максимум откликов.", GREEN),
]
cw = (SLIDE_W - 2 * MARGIN - 3 * 0.3) / 4
cy = 2.3
ch = 3.6
for i, (t, b, a) in enumerate(cards):
    add_card(s, MARGIN + i * (cw + 0.3), cy, cw, ch, t, b, accent=a, title_size=15, body_size=13)
add_footer(s, next_page())

# ======================================================================
# Slide 4 — Data
# ======================================================================
s = new_slide()
add_kicker(s, "Data Science", color=PRIMARY)
add_title(s, "X5 RetailHero: данные и постановка")
add_subtitle(s, "Реальный, публичный датасет — потоковая обработка 4.4 ГБ без загрузки в память целиком")

sw = (SLIDE_W - 2 * MARGIN - 3 * 0.3) / 4
sy = 2.3
sh = 2.0
add_stat_tile(s, MARGIN + 0 * (sw + 0.3), sy, sw, sh, "400 162", "клиента", "45.8 млн строк истории покупок до кампании", PRIMARY)
add_stat_tile(s, MARGIN + 1 * (sw + 0.3), sy, sw, sh, "200 039", "с известным откликом", "группы с коммуникацией и без — примерно поровну (~50/50)", AMBER)
add_stat_tile(s, MARGIN + 2 * (sw + 0.3), sy, sw, sh, "+3.3 п.п.", "средний эффект", "разница отклика между группами — точка отсчёта для всех моделей", GOLD)
add_stat_tile(s, MARGIN + 3 * (sw + 0.3), sy, sw, sh, "11.2%", "утечка устранена", "first_redeem_date после отсечки — цензурировано", GREEN)

add_card(s, MARGIN, sy + sh + 0.3, SLIDE_W - 2 * MARGIN, 1.65,
         "Находка: утечка в официальном эталонном решении",
         "Поле first_redeem_date простирается на 8 месяцев дальше окна истории покупок (отсечка 2019-03-18) — значит, "
         "оно частично описывает события ПОСЛЕ коммуникации. Официальное решение соревнования использует его «как есть», "
         "поэтому в этом проекте оно цензурировано.",
         accent=PRIMARY_DARK, title_size=15, body_size=13)
add_footer(s, next_page())

# ======================================================================
# Slide 5 — Classical models
# ======================================================================
s = new_slide()
add_kicker(s, "Data Science", color=PRIMARY)
add_title(s, "Классические uplift-модели")
add_subtitle(s, "Четыре подхода, один и тот же train/val/holdout сплит и одни и те же метрики")

grid = [
    ("0 · Random targeting", "Случайный порядок клиентов — не модель, а точка отсчёта: любая рабочая модель должна её обыграть.", "#8C8386"),
    ("1 · Response CatBoost", "Предсказывает P(Y=1|X) — вероятность покупки без учёта коммуникации, а значит, повторяет ту же ошибку, о которой шла речь на слайде 3.", PRIMARY),
    ("2 · T-Learner", "Обучает две отдельные CatBoost-модели на группах с коммуникацией и без, а прирост считает как разницу их прогнозов: τ(x) = μ₁(x) − μ₀(x).", AMBER),
    ("3 · X-Learner", "4-стадийная реализация (Kunzel et al., 2019) — не разность прогнозов, а отдельные модели эффекта с взвешиванием по вероятности коммуникации (propensity).", GREEN),
]
gw = (SLIDE_W - 2 * MARGIN - 0.3) / 2
gh = 1.75
gy = 2.3
for i, (t, b, a) in enumerate(grid):
    gx = MARGIN + (i % 2) * (gw + 0.3)
    gyy = gy + (i // 2) * (gh + 0.3)
    add_card(s, gx, gyy, gw, gh, t, b, accent=a, title_size=15.5, body_size=13.5)
add_footer(s, next_page())

# ======================================================================
# Slide 6 — Transformer architecture
# ======================================================================
s = new_slide()
add_kicker(s, "Development & Engineering", color=PRIMARY)
add_title(s, "Transformer: shared representation, two heads")
add_subtitle(s, "Обязательная модель проекта — работает на сырой последовательности чеков клиента")

boxes = ["Последовательность\nчеков клиента", "Эмбеддинги события,\nпозиции и времени", "Transformer-\nкодировщик (2 слоя)", "Представление\nклиента"]
bx, by, bw, bh, gap_ = 0.55, 2.15, 2.65, 1.0, 0.22
box_colors = [PRIMARY_DARK, PRIMARY_DARK, PRIMARY_DARK, GREEN]
for i, (text, col) in enumerate(zip(boxes, box_colors)):
    x = bx + i * (bw + gap_)
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(by), Inches(bw), Inches(bh))
    shp.fill.solid(); shp.fill.fore_color.rgb = rgb(col); shp.line.fill.background(); shp.shadow.inherit = False
    tf = shp.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = text
    run.font.size = Pt(13.5); run.font.bold = True; run.font.name = FONT; run.font.color.rgb = rgb(WHITE)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    if i < len(boxes) - 1:
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + bw + 0.02), Inches(by + bh / 2 - 0.11), Inches(gap_ - 0.04), Inches(0.22))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = rgb(GRAY_TEXT); arrow.line.fill.background(); arrow.shadow.inherit = False

head_y = 3.55
for i, (label, sub, col) in enumerate([("Без коммуникации", "μ₀(x)", AMBER), ("С коммуникацией", "μ₁(x)", GOLD)]):
    x = 4.75 + i * 3.0
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(head_y), Inches(2.6), Inches(0.85))
    shp.fill.solid(); shp.fill.fore_color.rgb = rgb(col); shp.line.fill.background(); shp.shadow.inherit = False
    tf = shp.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = f"{label}\n{sub}"
    run.font.size = Pt(14); run.font.bold = True; run.font.name = FONT; run.font.color.rgb = rgb(BG)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

down_arrow = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(bx + 3 * (bw + gap_) - gap_ / 2 + 0.15), Inches(by + bh + 0.03), Inches(0.22), Inches(0.26))
down_arrow.fill.solid(); down_arrow.fill.fore_color.rgb = rgb(GRAY_TEXT); down_arrow.line.fill.background(); down_arrow.shadow.inherit = False

result_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.75), Inches(4.75), Inches(2.6), Inches(0.7))
result_box.fill.solid(); result_box.fill.fore_color.rgb = rgb(CARD_BG); result_box.line.color.rgb = rgb(PRIMARY); result_box.line.width = Pt(1.25); result_box.shadow.inherit = False
tf = result_box.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
run = p.add_run(); run.text = "uplift = μ₁(x) − μ₀(x)"
run.font.size = Pt(15); run.font.bold = True; run.font.name = FONT; run.font.color.rgb = rgb(WHITE)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

add_card(s, MARGIN, 5.65, SLIDE_W - 2 * MARGIN, 1.05,
         "Единица последовательности — чек, не товарная позиция",
         "Сравнил распределения длин по трём уровням: на товарном истории оказались слишком длинными (p95=312, max=2513) — "
         "поэтому выбрал чек: он компактнее (p95=54) и при этом сохраняет реальную гранулярность похода в магазин.",
         accent=PRIMARY, title_size=14, body_size=12)
add_footer(s, next_page())

# ======================================================================
# Slide 7 — Results
# ======================================================================
s = new_slide()
add_kicker(s, "Data Science · Results", color=PRIMARY)
add_title(s, "Результаты на отложенной выборке (n = 40 011)")
add_text(s, MARGIN, 1.5, SLIDE_W - 2 * MARGIN, 0.9,
         "Единая метрика для всех моделей: src/uplift_metrics.py, 9 unit-тестов\n"
         "Qini/AUUC: 0 — модель не лучше случайного порядка; больше — точнее находит нужных клиентов; отрицательное — модель вредит",
         13.5, color=GRAY_TEXT, line_spacing=1.15)

table_models = ["x_learner_catboost", "t_learner_catboost", "transformer_two_head", "random_targeting", "response_catboost"]
rows = 1 + len(table_models)
cols = 4
tbl_left, tbl_top, tbl_w, tbl_h = MARGIN, 2.75, 6.0, 2.45
gframe = s.shapes.add_table(rows, cols, Inches(tbl_left), Inches(tbl_top), Inches(tbl_w), Inches(tbl_h))
table = gframe.table
headers = ["Модель", "Qini", "AUUC", "uplift@30%"]
for c, htext in enumerate(headers):
    cell = table.cell(0, c)
    cell.fill.solid(); cell.fill.fore_color.rgb = rgb(PRIMARY)
    cell.text = htext
    for p in cell.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.bold = True; run.font.size = Pt(13.5); run.font.name = FONT; run.font.color.rgb = rgb(WHITE)

for r, name in enumerate(table_models, start=1):
    m = metrics_by_model[name]
    vals = [MODEL_LABELS[name], f"{m['qini']:.4f}", f"{m['auuc']:.4f}", f"{m['uplift_at_30']:.4f}"]
    row_bg = CARD_BG if r % 2 else BG_ALT
    for c, v in enumerate(vals):
        cell = table.cell(r, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = rgb(row_bg)
        cell.text = v
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            for run in p.runs:
                run.font.size = Pt(12.5); run.font.name = FONT; run.font.color.rgb = rgb(WHITE)
                if name == "response_catboost" and c == 1:
                    run.font.color.rgb = rgb(MODEL_COLORS["response_catboost"]); run.font.bold = True
                if name == "x_learner_catboost":
                    run.font.bold = True
                    if c == 1:
                        run.font.color.rgb = rgb(PRIMARY)

add_image_fit(s, "reports/figures/presentation_qini_curves.png", tbl_left + tbl_w + 0.25, 2.75, 6.05, 4.15)
add_card(s, MARGIN, 5.3, tbl_w, 1.6,
         "Главный вывод",
         "Порядок X-Learner > T-Learner > Transformer > Random > Response одинаков по всем трём метрикам — не случайность "
         "одного показателя. Модель отклика значимо ХУЖЕ случайного распределения: весь ДИ Qini отрицательный.",
         accent=PRIMARY_DARK, title_size=13.5, body_size=11)
add_footer(s, next_page())

# ======================================================================
# Slide 8 — Business + A/B
# ======================================================================
s = new_slide()
add_kicker(s, "Product Thinking", color=PRIMARY)
add_title(s, "Бизнес-эффект (сценарий) и A/B-тест")
add_subtitle(s, "Офлайн Qini/AUUC — не доказательство реального бизнес-эффекта")

add_image_fit(s, "reports/figures/presentation_business_ev_curves.png", MARGIN, 2.2, 6.3, 4.55)
cards8 = [
    ("Сценарий, не факт", "Чтобы перевести ранжирование в деньги, я задал иллюстративные допущения — маржа 300₽, стоимость SMS 5₽. Это метод, а не факты X5 или Т-Банка.", AMBER),
    ("Лучший результат", "При этих допущениях наибольший ожидаемый эффект даёт X-Learner — на охвате примерно 57% аудитории.", GREEN),
    ("Перед внедрением", "Офлайн-метрики этого не доказывают, поэтому нужен онлайн A/B-эксперимент с защитными метриками и оценкой минимального эффекта.", PRIMARY),
]
cx = MARGIN + 6.3 + 0.3
cw8 = SLIDE_W - MARGIN - cx
cy8 = 2.2
ch8 = 1.45
for i, (t, b, a) in enumerate(cards8):
    add_card(s, cx, cy8 + i * (ch8 + 0.15), cw8, ch8, t, b, accent=a, title_size=14, body_size=12.5)
add_footer(s, next_page())

# ======================================================================
# Slide 9 — Conclusions
# ======================================================================
s = new_slide()
add_kicker(s, "Итог", color=GOLD)
add_title(s, "Выводы и ограничения")
add_subtitle(s, "Ранжирование моделей согласовано по всем метрикам, оба ключевых вывода подтверждены доверительными интервалами")

concl = [
    ("Ранжирование", "X-Learner > T-Learner > Transformer > Random > Response — согласовано по Qini, AUUC, uplift@30%.", PRIMARY),
    ("Отклик вреден", "Доверительный интервал Qini для модели отклика целиком отрицательный — значит, она не просто не помогает, а систематически хуже случайного распределения.", AMBER),
    ("Transformer уступает", "ДИ не пересекает ноль, значит, модель значимо лучше случайного распределения — но при этом CPU-бюджете обучения всё же уступает табличным моделям.", GOLD),
    ("Дальше", "Онлайн A/B-эксперимент, больший бюджет обучения Transformer, MLOps-пайплайн для продакшна.", GREEN),
]
cw9 = (SLIDE_W - 2 * MARGIN - 3 * 0.3) / 4
cy9 = 2.3
ch9 = 3.6
for i, (t, b, a) in enumerate(concl):
    add_card(s, MARGIN + i * (cw9 + 0.3), cy9, cw9, ch9, t, b, accent=a, title_size=14, body_size=12.5)
add_footer(s, next_page())

os.makedirs("presentation", exist_ok=True)
prs.save("presentation/presentation.pptx")
print("Saved presentation/presentation.pptx with", len(prs.slides._sldIdLst), "slides")
